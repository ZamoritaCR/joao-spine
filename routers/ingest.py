"""
JOAO Universal File Ingestion Pipeline
Accepts any number of files of any type.
Routes each file to the correct processor.
Returns structured results.
"""

import os
import io
import json
import tempfile
import mimetypes
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse

router = APIRouter()

# ─────────────────────────────────────────────
# File type routing
# ─────────────────────────────────────────────

AUDIO_EXTS = {'.mp3', '.wav', '.m4a', '.ogg', '.flac', '.aac', '.weba', '.opus'}
VIDEO_EXTS = {'.mp4', '.webm', '.mov', '.avi', '.mkv', '.m4v', '.wmv'}
DOC_EXTS   = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.md', '.rtf'}
DATA_EXTS  = {'.csv', '.xlsx', '.xls', '.json', '.parquet', '.tsv', '.jsonl'}
IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp', '.svg'}
CODE_EXTS  = {'.py', '.js', '.ts', '.html', '.css', '.sh', '.yaml', '.yml',
              '.toml', '.ini', '.env', '.sql', '.go', '.rs', '.java', '.cpp', '.c'}


def detect_type(filename: str, content_type: str = "") -> str:
    ext = Path(filename).suffix.lower()
    if ext in AUDIO_EXTS: return "audio"
    if ext in VIDEO_EXTS: return "video"
    if ext in DOC_EXTS:   return "document"
    if ext in DATA_EXTS:  return "data"
    if ext in IMAGE_EXTS: return "image"
    if ext in CODE_EXTS:  return "code"
    if "audio" in content_type: return "audio"
    if "video" in content_type: return "video"
    if "image" in content_type: return "image"
    return "text"


# ─────────────────────────────────────────────
# Processors
# ─────────────────────────────────────────────

async def process_audio_video(file_bytes: bytes, filename: str, file_type: str) -> dict:
    """Transcribe audio/video using OpenAI Whisper API."""
    import httpx
    from ..main import get_env  # noqa — pulls OPENAI_API_KEY from env

    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        # Fallback: just acknowledge, no transcription
        return {
            "type": file_type,
            "filename": filename,
            "status": "no_openai_key",
            "summary": "Audio/video received but no OpenAI key configured for Whisper transcription.",
            "transcript": None
        }

    with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        async with httpx.AsyncClient(timeout=120) as client:
            with open(tmp_path, "rb") as f:
                response = await client.post(
                    "https://api.openai.com/v1/audio/transcriptions",
                    headers={"Authorization": f"Bearer {api_key}"},
                    files={"file": (filename, f, "application/octet-stream")},
                    data={"model": "whisper-1", "response_format": "text"}
                )
        transcript = response.text if response.status_code == 200 else f"Whisper error: {response.text}"
    finally:
        os.unlink(tmp_path)

    # Summarize the transcript via Claude
    summary = await summarize_text(transcript[:8000], f"This is a transcript from {filename}. Summarize the key points concisely.")

    return {
        "type": file_type,
        "filename": filename,
        "status": "transcribed",
        "transcript": transcript,
        "summary": summary
    }


async def process_document(file_bytes: bytes, filename: str) -> dict:
    """Extract text from PDF, DOCX, PPTX, TXT, MD."""
    ext = Path(filename).suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext in {".docx", ".doc"}:
            import docx2txt
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                text = docx2txt.process(tmp_path)
            finally:
                os.unlink(tmp_path)

        elif ext in {".pptx", ".ppt"}:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append(f"Slide {i}: " + " | ".join(slide_text))
            text = "\n".join(slides)

        elif ext in {".txt", ".md", ".rtf"}:
            text = file_bytes.decode("utf-8", errors="replace")

    except Exception as e:
        text = f"[Extraction error: {str(e)}]"

    if not text.strip():
        text = "[No text could be extracted from this file]"

    summary = await summarize_text(text[:12000], f"Summarize this document ({filename}) — key points, main topics, any action items.")

    return {
        "type": "document",
        "filename": filename,
        "status": "extracted",
        "text": text[:5000] + ("..." if len(text) > 5000 else ""),
        "full_length": len(text),
        "summary": summary
    }


async def process_data_file(file_bytes: bytes, filename: str) -> dict:
    """Profile CSV, Excel, JSON, Parquet data files."""
    ext = Path(filename).suffix.lower()
    info = {}

    try:
        import pandas as pd

        if ext == ".csv":
            df = pd.read_csv(io.BytesIO(file_bytes))
        elif ext in {".xlsx", ".xls"}:
            df = pd.read_excel(io.BytesIO(file_bytes))
        elif ext == ".json":
            df = pd.read_json(io.BytesIO(file_bytes))
        elif ext == ".jsonl":
            df = pd.read_json(io.BytesIO(file_bytes), lines=True)
        elif ext == ".parquet":
            df = pd.read_parquet(io.BytesIO(file_bytes))
        elif ext == ".tsv":
            df = pd.read_csv(io.BytesIO(file_bytes), sep="\t")
        else:
            df = pd.read_csv(io.BytesIO(file_bytes))

        info = {
            "rows": int(df.shape[0]),
            "columns": int(df.shape[1]),
            "column_names": list(df.columns),
            "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
            "null_counts": {col: int(n) for col, n in df.isnull().sum().items()},
            "sample": df.head(5).to_dict(orient="records")
        }

        profile_text = f"""
Data file: {filename}
Rows: {info['rows']} | Columns: {info['columns']}
Columns: {', '.join(info['column_names'])}
Null counts: {json.dumps(info['null_counts'])}
Sample (first 5 rows): {json.dumps(info['sample'], default=str)}
        """.strip()

        summary = await summarize_text(profile_text, "Analyze this dataset and describe: what it contains, data quality issues, key patterns or insights.")

    except Exception as e:
        info = {"error": str(e)}
        summary = f"Could not parse data file: {str(e)}"

    return {
        "type": "data",
        "filename": filename,
        "status": "profiled",
        "profile": info,
        "summary": summary
    }


async def process_image(file_bytes: bytes, filename: str) -> dict:
    """Describe image using Claude vision."""
    import anthropic
    import base64

    ext = Path(filename).suffix.lower()
    media_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp"
    }
    media_type = media_map.get(ext, "image/jpeg")

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        return {
            "type": "image",
            "filename": filename,
            "status": "no_anthropic_key",
            "description": "Image received but no Anthropic API key configured."
        }

    try:
        client = anthropic.AsyncAnthropic(api_key=api_key)
        b64 = base64.standard_b64encode(file_bytes).decode("utf-8")

        message = await client.messages.create(
            model="claude-sonnet-4-5",
            max_tokens=1024,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {"type": "base64", "media_type": media_type, "data": b64}
                    },
                    {
                        "type": "text",
                        "text": "Describe this image in detail. What do you see? If it contains text, extract it. If it's a chart or diagram, explain what it shows."
                    }
                ]
            }]
        )
        description = message.content[0].text

    except Exception as e:
        description = f"Vision error: {str(e)}"

    return {
        "type": "image",
        "filename": filename,
        "status": "analyzed",
        "description": description
    }


async def process_code(file_bytes: bytes, filename: str) -> dict:
    """Analyze code files."""
    text = file_bytes.decode("utf-8", errors="replace")
    ext = Path(filename).suffix.lower()
    lang_map = {
        ".py": "Python", ".js": "JavaScript", ".ts": "TypeScript",
        ".html": "HTML", ".css": "CSS", ".sh": "Shell", ".sql": "SQL",
        ".go": "Go", ".rs": "Rust", ".java": "Java"
    }
    lang = lang_map.get(ext, "code")

    summary = await summarize_text(
        text[:8000],
        f"This is a {lang} file named {filename}. Explain what it does, its main functions/classes, and any notable patterns or issues."
    )

    return {
        "type": "code",
        "filename": filename,
        "language": lang,
        "status": "analyzed",
        "lines": len(text.splitlines()),
        "preview": text[:2000] + ("..." if len(text) > 2000 else ""),
        "summary": summary
    }


async def process_text(file_bytes: bytes, filename: str) -> dict:
    """Generic text file handler."""
    text = file_bytes.decode("utf-8", errors="replace")
    summary = await summarize_text(text[:8000], f"Summarize the contents of this file: {filename}")
    return {
        "type": "text",
        "filename": filename,
        "status": "read",
        "preview": text[:3000] + ("..." if len(text) > 3000 else ""),
        "summary": summary
    }


# ─────────────────────────────────────────────
# Claude summarizer (shared)
# ─────────────────────────────────────────────

async def summarize_text(text: str, instruction: str) -> str:
    """Send text to Claude for summarization/analysis."""
    import anthropic

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        return "[No Anthropic key — summary unavailable]"

    try:
        client = anthropic.AsyncAnthropic(api_key=api_key)
        message = await client.messages.create(
            model="claude-sonnet-4-5",
            max_tokens=2048,
            messages=[{
                "role": "user",
                "content": f"{instruction}\n\n---\n\n{text}"
            }]
        )
        return message.content[0].text
    except Exception as e:
        return f"[Summary error: {str(e)}]"


# ─────────────────────────────────────────────
# Route: POST /joao/ingest
# ─────────────────────────────────────────────

@router.post("/ingest")
async def ingest_files(
    files: List[UploadFile] = File(...),
    context: Optional[str] = Form(None)
):
    """
    Accept any number of files of any type.
    Auto-detect type and process each one.
    Return structured results for all files.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    results = []

    for upload in files:
        filename = upload.filename or "unnamed_file"
        content_type = upload.content_type or ""
        file_bytes = await upload.read()
        file_type = detect_type(filename, content_type)

        try:
            if file_type in ("audio", "video"):
                result = await process_audio_video(file_bytes, filename, file_type)
            elif file_type == "document":
                result = await process_document(file_bytes, filename)
            elif file_type == "data":
                result = await process_data_file(file_bytes, filename)
            elif file_type == "image":
                result = await process_image(file_bytes, filename)
            elif file_type == "code":
                result = await process_code(file_bytes, filename)
            else:
                result = await process_text(file_bytes, filename)

            result["size_bytes"] = len(file_bytes)

        except Exception as e:
            result = {
                "type": file_type,
                "filename": filename,
                "status": "error",
                "error": str(e),
                "size_bytes": len(file_bytes)
            }

        results.append(result)

    return JSONResponse({
        "files_processed": len(results),
        "context": context,
        "results": results
    })
